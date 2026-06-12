"""
Bluestock Mutual Fund Analytics Platform - PPTX Presentation Generator
Generates a professional 12-slide PowerPoint presentation deck summarizing
the project goals, data pipeline, performance metrics, and dashboard insights.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Set up paths relative to the script's directory
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)
REPORTS_DIR = os.path.join(PROJECT_ROOT, "reports")
CHARTS_DIR = os.path.join(REPORTS_DIR, "charts")
DASHBOARD_DIR = os.path.join(PROJECT_ROOT, "dashboard")
DAY3_CHARTS_DIR = os.path.join(REPORTS_DIR, "charts")

os.makedirs(REPORTS_DIR, exist_ok=True)
PPTX_OUTPUT_PATH = os.path.join(REPORTS_DIR, "Bluestock_MF_Presentation.pptx")

# Color definitions
NAVY = RGBColor(15, 23, 42)      # Primary Slate 900
BLUE = RGBColor(30, 58, 138)     # Secondary Blue 900
TEAL = RGBColor(13, 148, 136)    # Accent Teal 600
TEXT_DARK = RGBColor(51, 65, 85) # Slate 700
LIGHT_BG = RGBColor(248, 250, 252) # Slate 50

def add_header(slide, title_text):
    """Adds a standard header to a content slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

def add_bullet_points(slide, points, left, top, width, height, font_size=16):
    """Helper to add bullet points inside a textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, pt in enumerate(points):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_DARK
        p.level = 0
        p.space_after = Pt(8)

def format_dark_slide(slide, title_text, subtitle_text):
    """Sets a dark navy background and draws centered title/subtitle (for Cover/Ending)."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9.0), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.name = "Arial"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEAL
    p2.space_before = Pt(20)

def create_presentation():
    print("=" * 60)
    print("STARTING PRESENTATION GENERATION (pptx)")
    print("=" * 60)
    
    prs = Presentation()
    # Slide Dimensions: 10 x 7.5 inches (standard 4:3) or we can set it to Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # SLIDE 1: TITLE SLIDE (Dark Background)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    format_dark_slide(slide, "Mutual Fund Analytics Platform", "End-to-End Data Engineering & BI Dashboard | Bluestock Fintech")
    
    # Add metadata to title slide
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9.0), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Prepared by: Intern / Data Analyst\nDate: June 2026"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.space_after = Pt(4)
    
    # ----------------------------------------------------
    # SLIDE 2: PROBLEM & OBJECTIVE
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Problem Statement & Project Objectives")
    
    points_left = [
        "<b>Core Problems Addressed:</b>",
        "• <b>Data Fragmentation:</b> NAV, AUM, and SIP data scattered in TXT/PDF formats across AMFI & AMC portals.",
        "• <b>Performance Analysis Gap:</b> Retail investors struggle to evaluate risk-adjusted statistics (Sharpe, Sortino, Alpha, Beta).",
        "• <b>Lack of Benchmark Tracking:</b> Inability to verify index outperformance (Alpha & tracking error vs Nifty indices).",
        "• <b>Investor Behaviour Blind Spot:</b> Limited geographic and demographic visibility for targeted campaigns."
    ]
    add_bullet_points(slide, points_left, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    
    points_right = [
        "<b>Key Objectives & Outcomes:</b>",
        "• <b>O1:</b> Ingest daily NAV data for 40 mutual funds via public API pipelines.",
        "• <b>O2:</b> Build a normalized star schema SQLite relational database.",
        "• <b>O3:</b> Perform detailed EDA on NAV, AUM, and investor transactions.",
        "• <b>O4:</b> Compute performance scorecards and downside Value at Risk (VaR).",
        "• <b>O5:</b> Develop a 4-page interactive dashboard with filters."
    ]
    add_bullet_points(slide, points_right, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0))
    
    # ----------------------------------------------------
    # SLIDE 3: DATA SOURCES
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Data Sources & Project Inventory")
    
    points_left = [
        "<b>Primary Sources:</b>",
        "• <b>AMFI India:</b> Daily scheme NAVs & monthly aggregate notes.",
        "• <b>mfapi.in REST API:</b> Anchored historical NAV JSON feeds.",
        "• <b>NSE/BSE India:</b> Daily index closing prices (Nifty 50, Nifty 100).",
        "• <b>Synthesized Logs:</b> 32,000+ transaction logs modeled after real demographic distributions."
    ]
    add_bullet_points(slide, points_left, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    
    points_right = [
        "<b>Core Datasets Integrated:</b>",
        "1. <b>01_fund_master:</b> Scheme attributes (AMC, category, expense, manager).",
        "2. <b>02_nav_history:</b> Daily NAV data from Jan 2022 to May 2026 (~46,000 rows).",
        "3. <b>03_aum_by_fund_house:</b> Quarterly AUM figures for top 10 AMCs.",
        "4. <b>04_monthly_sip_inflows:</b> Month-wise SIP inflows & registrations.",
        "5. <b>08_investor_transactions:</b> Transactions for 5,000 investors across 12 states.",
        "6. <b>09_portfolio_holdings:</b> Stock holding weights & sector allocations."
    ]
    add_bullet_points(slide, points_right, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0))
    
    # ----------------------------------------------------
    # SLIDE 4: ARCHITECTURE
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "System Architecture & ETL Pipeline")
    
    points_left = [
        "<b>Data Pipeline Workflow:</b>",
        "• <b>1. Extract:</b> Python scripts fetch live NAVs from API (mfapi.in) and parse static flat files.",
        "• <b>2. Transform:</b> Gaps (holidays/weekends) resolved using forward-fill; dates standardized; columns aligned and validated.",
        "• <b>3. Load:</b> Schema created via schema.sql and populated into SQLite database (bluestock_mf.db).",
        "• <b>4. Analyze:</b> Calculations for CAGRs, Sharpe/Sortino ratios, OLS index regressions, VaR/CVaR, and sector HHI.",
        "• <b>5. Visualise:</b> Power BI imports data for multi-dimensional slicing."
    ]
    add_bullet_points(slide, points_left, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    
    points_right = [
        "<b>Relational Schema Highlights:</b>",
        "• <b>dim_fund:</b> Static master reference (PK: amfi_code).",
        "• <b>dim_date:</b> Date dimension (PK: date).",
        "• <b>fact_nav:</b> NAV facts (FKs: amfi_code, date).",
        "• <b>fact_transactions:</b> Investor transactions (FKs: amfi_code, date).",
        "• <b>fact_performance:</b> Quantitative performance facts.",
        "• <b>fact_portfolio:</b> Stock holdings allocations facts.",
        "• <b>Indexes:</b> Optimizations on fact_nav(amfi_code, date) and transaction keys."
    ]
    add_bullet_points(slide, points_right, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0))
    
    # ----------------------------------------------------
    # SLIDE 5: EDA HIGHLIGHTS (x1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "EDA Highlights I: Industry SIP Growth Trends")
    
    points_left = [
        "<b>Industry-wide Trends (2022-2025):</b>",
        "• <b>Assets Growth:</b> Industry AUM grew to Rs. 81 lakh crore in December 2025.",
        "• <b>SIP Expansion:</b> Active SIP accounts reached 9.35 crore, and monthly inflows peaked at Rs. 31,002 crore.",
        "• <b>Folio Inflows:</b> Folio counts increased from 13.26 crore to 26.12 crore, showing growth in retail investor interest.",
        "• <b>Category Preference:</b> Equity categories (Mid/Small Cap) led net inflows in FY 2024-25, followed by Index/ETF funds."
    ]
    add_bullet_points(slide, points_left, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    
    # Embed category inflow heatmap chart if available
    heatmap = os.path.join(DAY3_CHARTS_DIR, "06_category_inflow_heatmap.png")
    if os.path.exists(heatmap):
        slide.shapes.add_picture(heatmap, Inches(6.8), Inches(1.8), width=Inches(6.0), height=Inches(4.2))
        
    # ----------------------------------------------------
    # SLIDE 6: EDA HIGHLIGHTS (x2)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "EDA Highlights II: Investor Demographics & Geographics")
    
    points_left = [
        "<b>Investor Profiles (5,000 Sample):</b>",
        "• <b>Demographic Splits:</b> Age group 26-35 represents the largest investor base (38%), followed by 36-45 (29%).",
        "• <b>Geography Contribution:</b> T30 (Top 30 cities) accounts for 62% of investment volume. B30 (Beyond Top 30) is growing rapidly.",
        "• <b>Regional Leads:</b> Maharashtra, Gujarat, and Karnataka are the top states by transaction volume.",
        "• <b>Regular Contributions:</b> SIP transactions account for 78% of the transaction volume."
    ]
    add_bullet_points(slide, points_left, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    
    # Embed investor age pie chart
    pie = os.path.join(DAY3_CHARTS_DIR, "07_investor_age_pie.png")
    if os.path.exists(pie):
        slide.shapes.add_picture(pie, Inches(6.8), Inches(1.8), width=Inches(6.0), height=Inches(4.2))
        
    # ----------------------------------------------------
    # SLIDE 7: PERFORMANCE METRICS (x1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Performance Analytics I: Scorecard Rankings & CAGR")
    
    points_left = [
        "<b>Fund Scorecard Methodology:</b>",
        "• <b>Composite Scoring:</b> Calculated based on 3yr Return (30%), Sharpe (25%), Alpha (20%), Expense Ratio (15%, inverse), and Max Drawdown (10%, inverse).",
        "• <b>Top Performing Funds:</b>",
        "  1. <b>HDFC Top 100 Fund:</b> Ranked 1st in large-cap category.",
        "  2. <b>Kotak Emerging Equity:</b> Outperformed mid-cap peers.",
        "  3. <b>SBI Small Cap Fund:</b> Ranked 1st in small-cap category.",
        "• <b>Index Outperformance:</b> Active mid and small-cap managers generated significant alpha over their benchmarks."
    ]
    add_bullet_points(slide, points_left, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    
    # Embed benchmark comparison chart
    bench_comparison = os.path.join(CHARTS_DIR, "benchmark_comparison.png")
    if os.path.exists(bench_comparison):
        slide.shapes.add_picture(bench_comparison, Inches(6.8), Inches(1.8), width=Inches(6.0), height=Inches(4.2))
        
    # ----------------------------------------------------
    # SLIDE 8: PERFORMANCE METRICS (x2)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Performance Analytics II: Risk Metrics (Sharpe, VaR & HHI)")
    
    points_left = [
        "<b>Risk Metrics and Concentration Analysis:</b>",
        "• <b>Sharpe & Sortino:</b> Large-cap funds (HDFC Top 100, Sharpe = 1.06) showed higher risk-adjusted efficiency than small-cap funds.",
        "• <b>Value at Risk (95% VaR):</b> Small-cap funds showed the highest downside risk, with a 95% daily VaR of ~2.5%. Debt and liquid funds remained stable.",
        "• <b>Portfolio Concentration (HHI):</b> HHI values indicate that large-cap funds have higher sector concentration (HHI > 2,500). Mid and small-cap funds are more diversified (HHI ~1,500)."
    ]
    add_bullet_points(slide, points_left, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    
    # Embed HHI chart
    hhi = os.path.join(CHARTS_DIR, "sector_hhi_concentration.png")
    if os.path.exists(hhi):
        slide.shapes.add_picture(hhi, Inches(6.8), Inches(1.8), width=Inches(6.0), height=Inches(4.2))
        
    # ----------------------------------------------------
    # SLIDE 9: DASHBOARD SCREENSHOTS (x1)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Dashboard Walkthrough I: Industry Overview & Performance")
    
    p1 = os.path.join(DASHBOARD_DIR, "page_1.png")
    if os.path.exists(p1):
        slide.shapes.add_picture(p1, Inches(0.5), Inches(1.8), width=Inches(6.0), height=Inches(3.8))
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(6.0), Inches(1.2))
        txBox.text_frame.word_wrap = True
        txBox.text_frame.text = "<b>Page 1: Industry Overview</b>\nKPIs: Total AUM, monthly SIP inflows, and folio counts. Line charts show industry growth."
        txBox.text_frame.paragraphs[0].font.name = "Arial"
        txBox.text_frame.paragraphs[0].font.size = Pt(13)
        txBox.text_frame.paragraphs[0].font.color.rgb = NAVY
        
    p2 = os.path.join(DASHBOARD_DIR, "page_2.png")
    if os.path.exists(p2):
        slide.shapes.add_picture(p2, Inches(6.8), Inches(1.8), width=Inches(6.0), height=Inches(3.8))
        txBox = slide.shapes.add_textbox(Inches(6.8), Inches(5.8), Inches(6.0), Inches(1.2))
        txBox.text_frame.word_wrap = True
        txBox.text_frame.text = "<b>Page 2: Fund Performance</b>\nScatter plot of Risk vs Return, and interactive filters for AMC, category, and plan."
        txBox.text_frame.paragraphs[0].font.name = "Arial"
        txBox.text_frame.paragraphs[0].font.size = Pt(13)
        txBox.text_frame.paragraphs[0].font.color.rgb = NAVY
        
    # ----------------------------------------------------
    # SLIDE 10: DASHBOARD SCREENSHOTS (x2)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Dashboard Walkthrough II: Investor Analytics & Trends")
    
    p3 = os.path.join(DASHBOARD_DIR, "page_3.png")
    if os.path.exists(p3):
        slide.shapes.add_picture(p3, Inches(0.5), Inches(1.8), width=Inches(6.0), height=Inches(3.8))
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(6.0), Inches(1.2))
        txBox.text_frame.word_wrap = True
        txBox.text_frame.text = "<b>Page 3: Investor Analytics</b>\nGeographic bar charts showing transaction volumes by state, product split (SIP vs Lumpsum), and age distributions."
        txBox.text_frame.paragraphs[0].font.name = "Arial"
        txBox.text_frame.paragraphs[0].font.size = Pt(13)
        txBox.text_frame.paragraphs[0].font.color.rgb = NAVY
        
    p4 = os.path.join(DASHBOARD_DIR, "page_4.png")
    if os.path.exists(p4):
        slide.shapes.add_picture(p4, Inches(6.8), Inches(1.8), width=Inches(6.0), height=Inches(3.8))
        txBox = slide.shapes.add_textbox(Inches(6.8), Inches(5.8), Inches(6.0), Inches(1.2))
        txBox.text_frame.word_wrap = True
        txBox.text_frame.text = "<b>Page 4: SIP & Market Trends</b>\nVisualizes the relationship between SIP inflows and Nifty 50 index levels, and category flow heatmaps."
        txBox.text_frame.paragraphs[0].font.name = "Arial"
        txBox.text_frame.paragraphs[0].font.size = Pt(13)
        txBox.text_frame.paragraphs[0].font.color.rgb = NAVY
        
    # ----------------------------------------------------
    # SLIDE 11: KEY BUSINESS FINDINGS
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Key Business Findings & Operational Recommendations")
    
    points_left = [
        "<b>Core Strategic Takeaways:</b>",
        "• <b>SIP Continuity Risk:</b> ~97.8% of investors have average gaps between SIP dates exceeding 35 days, indicating irregular contributions. AMCs should implement auto-pay mandates and notifications.",
        "• <b>Sector Concentration:</b> Large-cap funds showed higher sector concentration (HHI > 2,500). Wealth managers should monitor portfolio concentration.",
        "• <b>Emerging Cohorts:</b> Newer cohorts (2025) are committing higher monthly ticket sizes despite smaller cohort volume.",
        "• <b>Automation:</b> Deploy the ETL pipeline as a scheduled job to automate reporting."
    ]
    add_bullet_points(slide, points_left, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    
    # ----------------------------------------------------
    # SLIDE 12: THANK YOU (Dark Background)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    format_dark_slide(slide, "Thank You", "Questions & Discussion | Bluestock Mutual Fund Analytics Platform")
    
    # Add email or closing note
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9.0), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Website: bluestock.in\nPrepared for Capstone Submission, June 2026"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)
    
    prs.save(PPTX_OUTPUT_PATH)
    print("=" * 60)
    print(f"PRESENTATION GENERATED AT {PPTX_OUTPUT_PATH}")
    print("=" * 60)

if __name__ == '__main__':
    create_presentation()
